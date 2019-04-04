
#%%
from pptx import Presentation
from pptx.util import Inches
import pathlib


#%%
image_file_path = pathlib.Path('formatted_images')
image_file_list = list(image_file_path.glob("*.jpg"))


#%%
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, World!"
subtitle.text = "Check out my l337 automation skillz!"

for image_file in image_file_list:
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    left = top = Inches(1)
    width = Inches(8.0)
    pic = slide.shapes.add_picture(str(image_file), left, top, width=width)

prs.save('test.pptx')


